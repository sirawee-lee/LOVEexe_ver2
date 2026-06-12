# -*- coding: utf-8 -*-
"""Generate the LOVE.EXE 5-minute presentation (editable .pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----- Theme -------------------------------------------------------------
BG        = RGBColor(0x14, 0x10, 0x1F)   # deep purple-black
PANEL     = RGBColor(0x21, 0x1A, 0x33)   # slightly lighter panel
ACCENT    = RGBColor(0xFF, 0x5C, 0x93)   # love pink/magenta
ACCENT2   = RGBColor(0x7A, 0x5C, 0xFF)   # electric violet
TEXT      = RGBColor(0xEC, 0xE7, 0xF4)   # near-white
MUTED     = RGBColor(0xA9, 0x9F, 0xC4)   # muted lavender
RED       = RGBColor(0xFF, 0x3B, 0x3B)   # screenshot-todo red
FILLHINT  = RGBColor(0x6E, 0xC8, 0xFF)   # blue for "fill in" placeholders

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=Pt(1), rounded=False, dash=None):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = line_w
        if dash:
            d = shp.line._get_or_add_ln()
            pd = d.makeelement(qn('a:prstDash'), {'val': dash})
            d.append(pd)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=Pt(4), line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
            r.font.name = 'Arial'
    return tb


def R(t, size, color=TEXT, bold=False, italic=False):
    return (t, size, color, bold, italic)


def header(s, kicker, title, num):
    box(s, Inches(0.55), Inches(0.55), Inches(0.12), Inches(0.85), fill=ACCENT)
    text(s, Inches(0.85), Inches(0.5), Inches(11.0), Inches(0.45),
         [[R(kicker, 14, ACCENT, True)]])
    text(s, Inches(0.83), Inches(0.82), Inches(11.4), Inches(0.7),
         [[R(title, 30, TEXT, True)]])
    text(s, Inches(12.3), Inches(0.55), Inches(0.7), Inches(0.4),
         [[R(num, 12, MUTED, True)]], align=PP_ALIGN.RIGHT)
    box(s, Inches(0.85), Inches(1.5), Inches(11.6), Pt(1.4), fill=PANEL)


def shot(s, x, y, w, h, caption):
    """Red-bordered screenshot placeholder with red instruction text."""
    box(s, x, y, w, h, fill=PANEL, line=RED, line_w=Pt(1.75),
        rounded=True, dash='dash')
    text(s, x + Inches(0.15), y, w - Inches(0.3), h,
         [[R("📷  SCREENSHOT NEEDED", 13, RED, True)],
          [R(caption, 11.5, RED, False, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(6))


def bullets(s, x, y, w, h, items, size=15, gap=Pt(9)):
    paras = []
    for it in items:
        paras.append([R("▸  ", size, ACCENT, True), R(it[0], size, TEXT, True)]
                      + ([R("  —  " + it[1], size - 1.5, MUTED)] if len(it) > 1 and it[1] else []))
    text(s, x, y, w, h, paras, space_after=gap, line_spacing=1.05)


# ========================================================================
# SLIDE 1 — TITLE
# ========================================================================
s = slide()
box(s, 0, 0, EMU_W, Inches(0.18), fill=ACCENT)
box(s, 0, Inches(7.32), EMU_W, Inches(0.18), fill=ACCENT2)
text(s, Inches(1.0), Inches(2.15), Inches(11.3), Inches(1.6),
     [[R("LOVE", 76, TEXT, True), R(".EXE", 76, ACCENT, True)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(3.55), Inches(11.3), Inches(0.6),
     [[R("A Campus Dating Sim Powered by a Gauntlet of Mini-Games", 19, MUTED, False, True)]],
     align=PP_ALIGN.CENTER)
box(s, Inches(5.17), Inches(4.35), Inches(3.0), Pt(2), fill=ACCENT)
text(s, Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.5),
     [[R("Built with Cocos Creator 2.4.8", 15, TEXT, True)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(5.25), Inches(11.3), Inches(1.0),
     [[R("Group [#]  ·  Final Project — Game Technology", 15, TEXT, True)],
      [R("[ Team members: name 1 · name 2 · name 3 · name 4 ]", 13.5, FILLHINT, False, True)]],
     align=PP_ALIGN.CENTER, space_after=Pt(8))
shot(s, Inches(4.42), Inches(0.55), Inches(4.5), Inches(1.45),
     "Title-screen key art / game logo")

# ========================================================================
# SLIDE 2 — CONCEPT
# ========================================================================
s = slide()
header(s, "01  ·  CONCEPT", "What Is LOVE.EXE?", "2")
text(s, Inches(0.85), Inches(1.85), Inches(7.0), Inches(0.9),
     [[R("“Survive campus life, win hearts, and dodge disaster — one mini-game at a time.”",
         18, ACCENT, True, True)]], line_spacing=1.1)
bullets(s, Inches(0.85), Inches(3.05), Inches(7.0), Inches(3.6), [
    ("Genre", "2D narrative dating-sim + mini-game collection"),
    ("Core loop", "explore campus, talk to NPCs, raise affinity"),
    ("Challenge", "story beats trigger skill-based mini-games"),
    ("Stakes", "a final boss stage and multiple branching endings"),
    ("Tone", "playful, romantic, a little chaotic"),
])
shot(s, Inches(8.15), Inches(1.95), Inches(4.4), Inches(4.6),
     "Main menu / title screen with START button")

# ========================================================================
# SLIDE 3 — GAME FLOW (10%)
# ========================================================================
s = slide()
header(s, "02  ·  GAME FLOW  (10%)", "How the Game Plays Out", "3")
flow = ["Main\nMenu", "Campus\nOverworld", "Story &\nNPC Events", "Mini-Games",
        "Boss Stage\n(5 games)", "Final\nChoice", "Multiple\nEndings"]
n = len(flow)
x0, y0 = Inches(0.85), Inches(2.05)
bw, bh, gap = Inches(1.46), Inches(1.1), Inches(0.26)
for i, label in enumerate(flow):
    bx = x0 + i * (bw + gap)
    col = ACCENT if i % 2 == 0 else ACCENT2
    b = box(s, bx, y0, bw, bh, fill=PANEL, line=col, line_w=Pt(1.75), rounded=True)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for j, ln in enumerate(label.split("\n")):
        pr = p if j == 0 else tf.add_paragraph()
        pr.alignment = PP_ALIGN.CENTER
        run = pr.add_run(); run.text = ln
        run.font.size = Pt(12.5); run.font.bold = True
        run.font.color.rgb = TEXT; run.font.name = 'Arial'
    if i < n - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                bx + bw + Inches(0.02), y0 + Inches(0.42),
                                gap - Inches(0.04), Inches(0.26))
        ar.fill.solid(); ar.fill.fore_color.rgb = MUTED
        ar.line.fill.background(); ar.shadow.inherit = False
text(s, Inches(0.85), Inches(3.55), Inches(11.6), Inches(0.9),
     [[R("Seamless scene management with animated fade in / fade out transitions between every stage.",
         15, MUTED, False, True)]])
shot(s, Inches(0.85), Inches(4.35), Inches(5.7), Inches(2.55),
     "Campus overworld — player walking the map")
shot(s, Inches(6.75), Inches(4.35), Inches(5.7), Inches(2.55),
     "A fade / transition moment between two scenes")

# ========================================================================
# SLIDE 4 — PHYSICS (13%)
# ========================================================================
s = slide()
header(s, "03  ·  CORE ENGINE TECH", "Gravity & Collision  (13%)", "4")
bullets(s, Inches(0.85), Inches(1.85), Inches(7.0), Inches(4.5), [
    ("Custom gravity", "velocity integration each frame: vy += g · dt"),
    ("Hand-rolled collision", "AABB rectangles + point-in-polygon zones"),
    ("Overworld blocking", "CollisionZone nodes gate player movement"),
    ("Per-game physics", "Flappy gravity, catcher hit-tests, dog pathfinding"),
    ("No engine shortcut", "all motion & collision written by us, not the built-in physics engine"),
], size=15.5, gap=Pt(11))
text(s, Inches(0.85), Inches(6.05), Inches(7.0), Inches(1.0),
     [[R("Files: ", 12.5, MUTED, True),
       R("FlappyGame.js · CollisionZone.js · MainGame.js", 12.5, FILLHINT, False, True)]])
shot(s, Inches(8.15), Inches(1.95), Inches(4.4), Inches(4.7),
     "Gameplay frame showing collision / gravity in action "
     "(e.g. Flappy bird vs. pipes, or player blocked by a wall)")

# ========================================================================
# SLIDE 5 — MINI-GAME VARIETY
# ========================================================================
s = slide()
header(s, "04  ·  TECH BREADTH", "One Engine, Many Mechanics", "5")
games = [
    ("Type Blast", "live typing + WPM scoring"),
    ("Rhythm (Osu)", "timing-based clicking"),
    ("Catcher", "mouse-tracked catching"),
    ("NiuPai Sim", "resource mgmt + enemy AI"),
    ("Flappy", "gravity & obstacle dodging"),
    ("Racing", "scrolling-lane dodging"),
    ("Thread Needle", "precision aiming"),
    ("Counter / Don't Press", "reaction & restraint"),
]
gx0, gy0 = Inches(0.85), Inches(1.9)
cw, ch, cgx, cgy = Inches(2.78), Inches(1.5), Inches(0.18), Inches(0.2)
for i, (name, mech) in enumerate(games):
    col, row = i % 4, i // 4
    cx = gx0 + col * (cw + cgx)
    cy = gy0 + row * (ch + cgy)
    box(s, cx, cy, cw, ch, fill=PANEL, line=ACCENT2, line_w=Pt(1.25), rounded=True)
    text(s, cx + Inches(0.18), cy + Inches(0.16), cw - Inches(0.36), ch - Inches(0.3),
         [[R(name, 15, ACCENT, True)], [R(mech, 12, MUTED)]], space_after=Pt(5))
shot(s, Inches(0.85), Inches(5.25), Inches(11.6), Inches(1.65),
     "A 2–3 panel collage of different mini-games "
     "(pick the most visually distinct screens)")

# ========================================================================
# SLIDE 6 — AUDIO / ANIMATION / JUICE
# ========================================================================
s = slide()
header(s, "05  ·  FEEL & POLISH", "Audio, Animation & Game Juice", "6")
stat = [("8", "BGM tracks", "one per scene/mood"),
        ("21", "sound effects", "clicks, hits, pickups, UI"),
        ("12", "animation clips", "idle · walk · jump · action"),
        ("5", "particle/juice FX", "burst · dust · shake · float-text")]
sx0 = Inches(0.85)
sw, sgap = Inches(2.78), Inches(0.18)
for i, (num, lbl, sub) in enumerate(stat):
    cx = sx0 + i * (sw + sgap)
    box(s, cx, Inches(1.95), sw, Inches(1.75), fill=PANEL, line=ACCENT, line_w=Pt(1.5), rounded=True)
    text(s, cx, Inches(2.05), sw, Inches(0.85),
         [[R(num, 40, ACCENT, True)]], align=PP_ALIGN.CENTER)
    text(s, cx, Inches(2.95), sw, Inches(0.7),
         [[R(lbl, 14, TEXT, True)], [R(sub, 11, MUTED)]],
         align=PP_ALIGN.CENTER, space_after=Pt(2))
bullets(s, Inches(0.85), Inches(4.05), Inches(7.1), Inches(2.7), [
    ("Sprite animations", "character idle / walk / jump / action states"),
    ("Particle systems", "cc.ParticleSystem dust + smell emitters"),
    ("Juice layer", "screen shake, floating score text, hover pops"),
    ("Transitions", "tweened fade in / out between scenes"),
], size=14.5, gap=Pt(9))
shot(s, Inches(8.2), Inches(4.05), Inches(4.35), Inches(2.75),
     "A moment with particles / animation / floating-score juice visible")

# ========================================================================
# SLIDE 7 — SYSTEMS & PERSISTENCE
# ========================================================================
s = slide()
header(s, "06  ·  SYSTEMS", "Data, Story & Persistence", "7")
bullets(s, Inches(0.85), Inches(1.9), Inches(7.0), Inches(4.6), [
    ("Leaderboard", "top-10 scores saved in localStorage"),
    ("Branching dialogue", "choices tracked via a StoryState manager"),
    ("Affinity system", "relationship values steer the story"),
    ("Multiple endings", "outcome decided by player choices"),
    ("Pause menu", "ESC to pause / resume anytime"),
    ("Gamepad support", "playable with a controller"),
    ("Version control", "Git used throughout for team collaboration"),
], size=15, gap=Pt(10))
shot(s, Inches(8.15), Inches(1.95), Inches(4.4), Inches(2.25),
     "Leaderboard screen (top-10 scores)")
shot(s, Inches(8.15), Inches(4.35), Inches(4.4), Inches(2.25),
     "Dialogue / branching choice moment")

# ========================================================================
# SLIDE 8 — HIGHLIGHTS / BONUS
# ========================================================================
s = slide()
header(s, "07  ·  HIGHLIGHTS", "What We're Proud Of", "8")
bullets(s, Inches(0.85), Inches(1.95), Inches(7.0), Inches(4.4), [
    ("A cohesive story", "narrative threads every mini-game together"),
    ("Multiple endings", "distinct romance routes + a hidden easter egg"),
    ("Mechanic variety", "9 mini-games, each a different skill"),
    ("Custom UI", "interface drawn programmatically with cc.Graphics"),
    ("Team workflow", "Git branches & merges across the team"),
], size=15.5, gap=Pt(12))
shot(s, Inches(8.15), Inches(1.95), Inches(4.4), Inches(4.6),
     "An ending screen / favorite key moment")

# ========================================================================
# SLIDE 9 — THANK YOU
# ========================================================================
s = slide()
box(s, 0, 0, EMU_W, Inches(0.18), fill=ACCENT2)
box(s, 0, Inches(7.32), EMU_W, Inches(0.18), fill=ACCENT)
text(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.2),
     [[R("Thank You", 60, TEXT, True), R("  ♥", 60, ACCENT, True)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.6),
     [[R("Questions?  —  LOVE.EXE, built with Cocos Creator 2.4.8", 18, MUTED, False, True)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.5),
     [[R("[ Group # · team member names · repo link ]", 14, FILLHINT, False, True)]],
     align=PP_ALIGN.CENTER)

out = "/Users/admin/Desktop/LoveEXE_final/LOVE_EXE_Presentation.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
